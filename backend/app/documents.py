import io

from docx import Document as DocxDocument
from fastapi import APIRouter, Header, HTTPException
from pptx import Presentation
from pypdf import PdfReader

from .supabase_client import get_admin_client

router = APIRouter()

CHUNK_SIZE = 1200
CHUNK_OVERLAP = 150


def extract_text(file_bytes: bytes, filename: str) -> str:
    lower = filename.lower()
    if lower.endswith(".pdf"):
        reader = PdfReader(io.BytesIO(file_bytes))
        return "\n\n".join(page.extract_text() or "" for page in reader.pages)
    if lower.endswith(".pptx"):
        prs = Presentation(io.BytesIO(file_bytes))
        slides_text = []
        for slide in prs.slides:
            texts = [shape.text for shape in slide.shapes if shape.has_text_frame]
            slides_text.append("\n".join(texts))
        return "\n\n".join(slides_text)
    if lower.endswith(".docx"):
        doc = DocxDocument(io.BytesIO(file_bytes))
        return "\n\n".join(p.text for p in doc.paragraphs)
    if lower.endswith(".txt"):
        return file_bytes.decode("utf-8", errors="ignore")
    raise ValueError(f"Unsupported file type: {filename}")


def chunk_text(
    text: str, chunk_size: int = CHUNK_SIZE, overlap: int = CHUNK_OVERLAP
) -> list[str]:
    paragraphs = [p.strip() for p in text.split("\n\n") if p.strip()]
    chunks: list[str] = []
    current = ""
    for para in paragraphs:
        if len(current) + len(para) + 2 <= chunk_size:
            current = f"{current}\n\n{para}" if current else para
            continue

        if current:
            chunks.append(current)

        if len(para) > chunk_size:
            for i in range(0, len(para), chunk_size - overlap):
                chunks.append(para[i : i + chunk_size])
            current = ""
        else:
            current = para

    if current:
        chunks.append(current)
    return chunks


def get_user_id(authorization: str | None) -> str:
    if not authorization or not authorization.startswith("Bearer "):
        raise HTTPException(status_code=401, detail="Missing bearer token")
    token = authorization.removeprefix("Bearer ")
    admin = get_admin_client()
    result = admin.auth.get_user(token)
    if not result or not result.user:
        raise HTTPException(status_code=401, detail="Invalid token")
    return result.user.id


@router.delete("/documents/{document_id}")
def delete_document(
    document_id: str, authorization: str | None = Header(default=None)
):
    user_id = get_user_id(authorization)
    admin = get_admin_client()

    doc_result = (
        admin.table("documents")
        .select("*")
        .eq("id", document_id)
        .maybe_single()
        .execute()
    )
    document = doc_result.data if doc_result else None
    if not document:
        raise HTTPException(status_code=404, detail="Document not found")
    if document["user_id"] != user_id:
        raise HTTPException(status_code=403, detail="Not your document")

    concept_ids = [
        c["id"]
        for c in (
            admin.table("concepts")
            .select("id")
            .eq("document_id", document_id)
            .execute()
            .data
            or []
        )
    ]

    # Children before parents, mirroring DEMO_TABLES_IN_DELETE_ORDER in demo.py.
    if concept_ids:
        admin.table("quiz_responses").delete().in_("concept_id", concept_ids).execute()
        admin.table("concept_mastery").delete().in_(
            "concept_id", concept_ids
        ).execute()
    admin.table("quiz_questions").delete().eq("document_id", document_id).execute()
    admin.table("flashcards").delete().eq("document_id", document_id).execute()
    admin.table("concepts").delete().eq("document_id", document_id).execute()
    admin.table("document_chunks").delete().eq("document_id", document_id).execute()

    admin.storage.from_("study-materials").remove([document["storage_path"]])
    admin.table("documents").delete().eq("id", document_id).execute()

    return {"status": "deleted"}


@router.post("/documents/{document_id}/process")
def process_document(
    document_id: str, authorization: str | None = Header(default=None)
):
    user_id = get_user_id(authorization)
    admin = get_admin_client()

    doc_result = (
        admin.table("documents")
        .select("*")
        .eq("id", document_id)
        .maybe_single()
        .execute()
    )
    document = doc_result.data if doc_result else None
    if not document:
        raise HTTPException(status_code=404, detail="Document not found")
    if document["user_id"] != user_id:
        raise HTTPException(status_code=403, detail="Not your document")

    admin.table("documents").update({"status": "processing"}).eq(
        "id", document_id
    ).execute()

    try:
        file_bytes = admin.storage.from_("study-materials").download(
            document["storage_path"]
        )
        text = extract_text(file_bytes, document["filename"])
        if not text.strip():
            raise ValueError("No extractable text found in this file.")
        chunks = chunk_text(text)

        admin.table("document_chunks").delete().eq(
            "document_id", document_id
        ).execute()
        rows = [
            {
                "document_id": document_id,
                "user_id": user_id,
                "chunk_index": i,
                "content": chunk,
            }
            for i, chunk in enumerate(chunks)
        ]
        if rows:
            admin.table("document_chunks").insert(rows).execute()

        admin.table("documents").update(
            {"status": "processed", "error_message": None}
        ).eq("id", document_id).execute()
        return {"status": "processed", "chunk_count": len(chunks)}
    except Exception as exc:
        admin.table("documents").update(
            {"status": "error", "error_message": str(exc)}
        ).eq("id", document_id).execute()
        raise HTTPException(status_code=500, detail=str(exc))
